import os
import re
import logging
from io import BytesIO
from telegram import Update
from telegram.ext import Application, CommandHandler, MessageHandler, filters, ContextTypes, ConversationHandler
from bs4 import BeautifulSoup
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

# Enable logging
logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    level=logging.INFO
)
logger = logging.getLogger(__name__)

# Conversation states
WAITING_LINK, WAITING_IMAGE = range(2)

class DuaLine:
    def __init__(self, arabic, translation):
        self.arabic = arabic.strip()
        self.translation = translation.strip()

def scrape_dua(url):
    """Scrape dua content from duas.org"""
    try:
        response = requests.get(url, timeout=30)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'html.parser')
        
        dua_lines = []
        
        # Find all Arabic divs
        arabic_divs = soup.find_all('div', class_='Ara')
        translation_divs = soup.find_all('div', class_='Tra')
        
        # Match Arabic with English translations
        for i in range(min(len(arabic_divs), len(translation_divs))):
            arabic_text = arabic_divs[i].get_text(strip=True)
            translation_text = translation_divs[i].get_text(strip=True)
            
            if arabic_text and translation_text:
                dua_lines.append(DuaLine(arabic_text, translation_text))
        
        # Get title
        title_elem = soup.find('ptitle')
        title = title_elem.get_text(strip=True) if title_elem else "Dua"
        
        return dua_lines, title
    
    except Exception as e:
        logger.error(f"Error scraping dua: {e}")
        raise

def create_pptx_single_line(dua_lines, title, bg_image_path):
    """Create PPTX with one line per slide (Arabic left, English right)"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    
    for line in dua_lines:
        slide = prs.slides.add_slide(blank_layout)
        
        # Add background image
        if bg_image_path and os.path.exists(bg_image_path):
            slide.shapes.add_picture(
                bg_image_path, 
                0, 0, 
                width=prs.slide_width, 
                height=prs.slide_height
            )
        
        # Add semi-transparent overlay for better text readability
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            0, 0,
            prs.slide_width,
            prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
        shape.fill.transparency = 0.3
        shape.line.fill.background()
        
        # Arabic text (right side)
        arabic_box = slide.shapes.add_textbox(
            Inches(5.2), 
            Inches(2.5), 
            Inches(4.5), 
            Inches(2.5)
        )
        arabic_frame = arabic_box.text_frame
        arabic_frame.word_wrap = True
        arabic_frame.vertical_anchor = 1  # Middle
        
        p = arabic_frame.paragraphs[0]
        p.text = line.arabic
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = 'Arial'
        
        # English text (left side)
        english_box = slide.shapes.add_textbox(
            Inches(0.3), 
            Inches(2.5), 
            Inches(4.5), 
            Inches(2.5)
        )
        english_frame = english_box.text_frame
        english_frame.word_wrap = True
        english_frame.vertical_anchor = 1  # Middle
        
        p = english_frame.paragraphs[0]
        p.text = line.translation
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = 'Arial'
    
    return prs

def create_pptx_three_lines(dua_lines, title, bg_image_path):
    """Create PPTX with three lines per slide (Arabic right, English left)"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Group lines into sets of 3
    for i in range(0, len(dua_lines), 3):
        slide = prs.slides.add_slide(blank_layout)
        
        # Add background image
        if bg_image_path and os.path.exists(bg_image_path):
            slide.shapes.add_picture(
                bg_image_path, 
                0, 0, 
                width=prs.slide_width, 
                height=prs.slide_height
            )
        
        # Add semi-transparent overlay
        shape = slide.shapes.add_shape(
            1,
            0, 0,
            prs.slide_width,
            prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
        shape.fill.transparency = 0.3
        shape.line.fill.background()
        
        # Get up to 3 lines for this slide
        current_lines = dua_lines[i:i+3]
        
        # Arabic text box (right side)
        arabic_box = slide.shapes.add_textbox(
            Inches(5.2), 
            Inches(1), 
            Inches(4.5), 
            Inches(5.5)
        )
        arabic_frame = arabic_box.text_frame
        arabic_frame.word_wrap = True
        arabic_frame.vertical_anchor = 1  # Middle
        
        for idx, line in enumerate(current_lines):
            if idx > 0:
                arabic_frame.add_paragraph()
            p = arabic_frame.paragraphs[idx]
            p.text = line.arabic
            p.alignment = PP_ALIGN.RIGHT
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = 'Arial'
            p.space_after = Pt(20)
        
        # English text box (left side)
        english_box = slide.shapes.add_textbox(
            Inches(0.3), 
            Inches(1), 
            Inches(4.5), 
            Inches(5.5)
        )
        english_frame = english_box.text_frame
        english_frame.word_wrap = True
        english_frame.vertical_anchor = 1  # Middle
        
        for idx, line in enumerate(current_lines):
            if idx > 0:
                english_frame.add_paragraph()
            p = english_frame.paragraphs[idx]
            p.text = line.translation
            p.alignment = PP_ALIGN.LEFT
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = 'Arial'
            p.space_after = Pt(20)
    
    return prs

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Start command handler"""
    await update.message.reply_text(
        "🕌 Welcome to the Duas.org to PPTX Bot!\n\n"
        "Send me a link from duas.org (e.g., https://www.duas.org/kumayl.html)\n\n"
        "Use /cancel to stop the process at any time."
    )
    return WAITING_LINK

async def receive_link(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle the duas.org link"""
    url = update.message.text.strip()
    
    # Validate URL
    if not url.startswith('https://www.duas.org/') and not url.startswith('http://www.duas.org/'):
        await update.message.reply_text(
            "❌ Please send a valid duas.org link.\n"
            "Example: https://www.duas.org/kumayl.html"
        )
        return WAITING_LINK
    
    await update.message.reply_text("🔄 Scraping dua content...")
    
    try:
        dua_lines, title = scrape_dua(url)
        
        if not dua_lines:
            await update.message.reply_text(
                "❌ Could not extract dua content from this page. "
                "Please make sure the link contains a valid dua."
            )
            return WAITING_LINK
        
        # Store in context
        context.user_data['dua_lines'] = dua_lines
        context.user_data['title'] = title
        
        await update.message.reply_text(
            f"✅ Found {len(dua_lines)} lines from: {title}\n\n"
            "📷 Now send me a background image for the slides."
        )
        return WAITING_IMAGE
        
    except Exception as e:
        logger.error(f"Error processing link: {e}")
        await update.message.reply_text(
            f"❌ Error scraping the dua: {str(e)}\n"
            "Please try again with a different link."
        )
        return WAITING_LINK

async def receive_image(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle the background image and generate PPTXs"""
    if not update.message.photo:
        await update.message.reply_text("❌ Please send an image file.")
        return WAITING_IMAGE
    
    await update.message.reply_text("🔄 Processing image and generating presentations...")
    
    try:
        # Download image
        photo = update.message.photo[-1]  # Get highest resolution
        file = await context.bot.get_file(photo.file_id)
        
        # Save image temporarily
        image_path = f"temp_bg_{update.effective_user.id}.jpg"
        await file.download_to_drive(image_path)
        
        # Resize image if needed
        img = Image.open(image_path)
        img = img.resize((1280, 960), Image.Resampling.LANCZOS)
        img.save(image_path, quality=95)
        
        # Get stored data
        dua_lines = context.user_data.get('dua_lines', [])
        title = context.user_data.get('title', 'Dua')
        
        # Generate single-line PPTX
        prs_single = create_pptx_single_line(dua_lines, title, image_path)
        single_file = BytesIO()
        prs_single.save(single_file)
        single_file.seek(0)
        
        # Generate three-line PPTX
        prs_three = create_pptx_three_lines(dua_lines, title, image_path)
        three_file = BytesIO()
        prs_three.save(three_file)
        three_file.seek(0)
        
        # Send files
        await update.message.reply_document(
            document=single_file,
            filename=f"{title.replace(' ', '_')}_single_line.pptx",
            caption="📄 Single line per slide version"
        )
        
        await update.message.reply_document(
            document=three_file,
            filename=f"{title.replace(' ', '_')}_three_lines.pptx",
            caption="📄 Three lines per slide version"
        )
        
        # Cleanup
        if os.path.exists(image_path):
            os.remove(image_path)
        
        context.user_data.clear()
        
        await update.message.reply_text(
            "✅ Done! Your presentations are ready.\n\n"
            "Send another duas.org link to create more presentations, "
            "or use /cancel to stop."
        )
        
        return WAITING_LINK
        
    except Exception as e:
        logger.error(f"Error generating PPTX: {e}")
        await update.message.reply_text(
            f"❌ Error generating presentations: {str(e)}\n"
            "Please try again."
        )
        return WAITING_IMAGE

async def cancel(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Cancel the conversation"""
    context.user_data.clear()
    await update.message.reply_text(
        "❌ Operation cancelled. Use /start to begin again."
    )
    return ConversationHandler.END

def main():
    """Start the bot"""
    TOKEN = os.environ.get('TELEGRAM_BOT_TOKEN')
    
    if not TOKEN:
        raise ValueError("TELEGRAM_BOT_TOKEN environment variable not set!")
    
    # Create application
    application = Application.builder().token(TOKEN).build()
    
    # Conversation handler
    conv_handler = ConversationHandler(
        entry_points=[CommandHandler('start', start)],
        states={
            WAITING_LINK: [MessageHandler(filters.TEXT & ~filters.COMMAND, receive_link)],
            WAITING_IMAGE: [MessageHandler(filters.PHOTO, receive_image)],
        },
        fallbacks=[CommandHandler('cancel', cancel)],
    )
    
    application.add_handler(conv_handler)
    
    # Start bot with webhook for Render
    PORT = int(os.environ.get('PORT', 8443))
    WEBHOOK_URL = os.environ.get('WEBHOOK_URL')  # e.g., https://your-app.onrender.com
    
    if WEBHOOK_URL:
        logger.info(f"Starting webhook on port {PORT}")
        application.run_webhook(
            listen="0.0.0.0",
            port=PORT,
            webhook_url=f"{WEBHOOK_URL}/{TOKEN}",
            url_path=TOKEN
        )
    else:
        logger.info("Starting polling mode")
        application.run_polling()

if __name__ == '__main__':
    main()
